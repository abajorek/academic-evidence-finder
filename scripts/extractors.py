import os
from pdfminer.high_level import extract_text as pdf_extract_text
import docx2txt
from pptx import Presentation

def read_txt(path:str)->str:
    with open(path, "r", errors="ignore") as f:
        return f.read()

def read_pdf(path:str)->str:
    try:
        return pdf_extract_text(path) or ""
    except Exception:
        return ""

def read_docx(path:str)->str:
    try:
        return docx2txt.process(path) or ""
    except Exception:
        return ""

def read_pptx(path:str)->str:
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception:
        return ""

READERS = {
    ".txt": read_txt,
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
}

def extract_text(path:str)->str:
    _, ext = os.path.splitext(path.lower())
    reader = READERS.get(ext)
    if reader:
        return reader(path)
    return ""
